import time
from flask import Flask
from pptx import Presentation
from pptx.util import Cm, Pt
from pdfreader import PDFDocument
import io
import tesserocr
from PIL import Image


app = Flask(__name__)

@app.route('/time')
def get_time():
    return {'time' : time.time()}


@app.route('/ppt')
def ppt():
    prs = Presentation('/Users/bhavyameghnani/Desktop/IVI/misc/ivi.pptx')
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    print(text_runs)
    return "Check Console"


@app.route('/ppt_font')
def ppt_font():
    prs = Presentation('/Users/bhavyameghnani/Desktop/IVI/misc/ivi.pptx')

    # text_frame = shape.text_frame
    # paragraph = text_frame.paragraphs[0]

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font = run.font
                    try:
                        font_size = font.size
                        print(font_size)
                        font_bold = font.bold
                        print(font_bold)
                        font_name = font.name
                        print(font_name)
                        # color = font.color.rgb
                        # print(color)
                    except:
                        pass
    return "Check Console"



@app.route('/pdf')
def pdf():
    fd = open('/Users/bhavyameghnani/Desktop/IVI/misc/ivi.pdf', "rb")
    doc = PDFDocument(fd)
    page = next(doc.pages())
    KEY = sorted(page.Resources.Font.keys())
    print(KEY)
    for i in range(len(KEY)):
        font = page.Resources.Font[KEY[i]]
        print("FONT SUBTYPE "+font.Subtype)
        print("BASE FONT "+font.BaseFont)
        print("FONT ENCODING "+font.Encoding)
    

    return "Check Console"

@app.route('/img')
def img():
    with tesserocr.PyTessBaseAPI() as api:
        image = Image.open("/Users/bhavyameghnani/Desktop/IVI/misc/ivi.png")
        api.SetImage(image)
        api.Recognize()  # required to get result from the next line
        iterator = api.GetIterator()
        print(iterator.WordFontAttributes())
    return "Check Console"

app.run(port=5000, debug=True)


# cd flask
# source bin/activate